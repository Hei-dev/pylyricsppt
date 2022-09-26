#Temp. patch for Python 3.10 python-pptx
import collections 
import collections.abc
#Import required library
import pptx #https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
from pptx.util import Inches as ppt_inch
from pptx.util import Cm as ppt_cm
from pptx.util import Pt as ppt_pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

import PIL

import fnmatch
from collections import defaultdict
import json
from pptx.enum.dml import MSO_FILL
from os import path
import colorsys

prs = pptx.Presentation()
prs.slide_width = pptx.util.Inches(16)
prs.slide_height = pptx.util.Inches(9)

#image analysing
cacheImageData = {}

#Lyrics-related variables
lyrics = []

#Read the lyrics
lySlide = open("source/lyrics.txt","r",encoding="utf-8")
lyricsRaw = lySlide.read()
lySlide.close()

# Read and loads font preferences
pref_f = open("source/font.json","r",encoding="utf-8")
pref = json.load(pref_f)
pref_f.close()

ptTable = pref["size"]
max_line_len = pref["max_line_length"]

def addLyricsSlide(lyric,textType,imgPath):
    # Check if imga eneeds to be analysed
    if (picname not in cacheImageData) and (picname!=""):
        analysePic(picname)

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    title = slide.shapes.title

    title.width = ppt_inch(16)
    title.height = ppt_inch(1.45)

    title.top = int((prs.slide_height - title.height) / 2)
    title.left = ppt_inch(0)

    lyric_nopun = replaceAllPuncuation(lyric)

    if len(lyric_nopun)>=max_line_len:
        # if the lyric would auto new line (i.e. too many chars)
        lyric_multiline = lyric_nopun.split(" ")
        p = title.text_frame.paragraphs[0]
        p.text = lyric_multiline[0]
        for ly_l in lyric_multiline[1:]:
            p = title.text_frame.add_paragraph()
            p.text = ly_l
            titleFont = title.text_frame.paragraphs[lyric_multiline.index(ly_l)].font
            titleFont.size = ppt_pt(ptTable[textType])
            titleFont.name = pref["typeface"]

            if picname != "":
                titleFont.color.rgb = cacheImageData[imgPath]["textColor"]
    else:
        title.text = lyric_nopun

    titleLayout = title.text_frame.paragraphs[0]
    titleLayout.alignment = PP_ALIGN.CENTER
    titleFont = title.text_frame.paragraphs[0].font
    #titleFont2 = title.text_frame.paragraphs[0].font
    titleFont.size = ppt_pt(ptTable[textType])
    titleFont.name = pref["typeface"]
    if picname != "":
        titleFont.color.rgb = cacheImageData[imgPath]["textColor"]

    #TODO Make this background image at some point, not a shape
    if imgPath!="":
        bg = slide.shapes.add_picture(imgPath, 0, 0, prs.slide_width, prs.slide_height) 
        slide.shapes._spTree.insert(2, bg._element)

#print(lyricsRaw.splitlines())

#Process the lyrics: Categorizze them
songLy = {}
curTitle = ""
curVerse = ""
curVerseLy = []

def analysePic(picPath):
    #f_p = open(picPath,encoding="utf-16")
    picbg = PIL.Image.open(picPath).convert("HSV")
    cacheImageData[picPath] = {}
    #f_p.close()

    pic_w,pic_h = picbg.size
    pic_px = picbg.load()
    avg_v_middle = 0
    mid_count = 0
    for x in range(0,pic_w):
        for y in range(0,pic_h):
            #print(pic_px[x,y])
            #TODO analyse the min max of each H, S, V (analyse for whole pic and the middle (1/3) part)
            if y>=pic_h/3 and y<=pic_h/3*2:
                avg_v_middle += pic_px[x,y][2]
                mid_count += 1
    
    avg_v_middle = avg_v_middle / mid_count
    if avg_v_middle < 127:
        cacheImageData[picPath]["textColor"] = RGBColor(0x00,0x00,0x00)
    else:
        cacheImageData[picPath]["textColor"] = RGBColor(0xFF,0xFF,0xFF)

#analysePic("source/高唱入雲/V.jpg")

def appendToSongLy():
    global curVerseLy
    if curVerse!="" and curVerseLy!=[]:
        songLy[curTitle][curVerse] = curVerseLy
        curVerseLy = []

def replaceAllPuncuation(text):
    return text.replace("，"," ").replace("。"," ").replace("："," ").replace("."," ").replace(","," ").replace(":"," ").replace(";"," ")

for line in lyricsRaw.splitlines():
    #print(curVerse,line)
    if fnmatch.filter([line],"V?:") or fnmatch.filter([line],"V:") or fnmatch.filter([line],"C?:") or fnmatch.filter([line],"C:"):
        #print("Change",curTitle,curVerse)
        appendToSongLy()
        curVerse = line
    elif fnmatch.filter([line],"B?:") or fnmatch.filter([line],"B:"):
        appendToSongLy()
        curVerse = line
    elif "T:" in line:
        appendToSongLy()
        curTitle = line[2:]
        songLy[curTitle] = {}
    #elif line=="":
    #    curVerse = line
    else:
        if line!="":
            curVerseLy.append(line)
#Append the last section
appendToSongLy()

# Paste the lyrics to slides
for k_title in songLy.keys():
    if path.exists("source/"+k_title):
        if path.exists("source/"+k_title+"/T.jpg"):
            picname = "source/"+k_title+"/T.jpg"
        elif path.exists("source/"+k_title+"/T.jpeg"):
            picname = "source/"+k_title+"/T.jpeg"
        else:
            picname = ""
    else:
        picname = ""
    
    addLyricsSlide(k_title,"Title",picname)
    for k_verse_r in songLy[k_title].keys():
        k_verse = k_verse_r.replace(":","")
        path_to_pics = "source/"+k_title+"/"+k_verse.replace(":","")
        print("Loading background",path_to_pics)
        if path.exists("source/"+k_title):
            if path.exists(path_to_pics+".jpg"):
                picname = path_to_pics+".jpg"
            elif path.exists(path_to_pics+".jpeg"):
                picname = path_to_pics+".jpeg"
            else:
                picname = ""
        else:
            picname = ""

        for lyr in songLy[k_title][k_verse+":"]:
            addLyricsSlide(lyr,"Lyrics",picname)

print(songLy)



prs.save('test.pptx')